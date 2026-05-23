from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\angpa\myProjects\Daily_Work\AI_Tech_Review")
TOPIC = "2026-04-17_ai-updates-weekly"
TOPIC_DIR = ROOT / TOPIC
TEMPLATE = TOPIC_DIR / "skywork_inputs" / "LGD_Template.pptx"
OUTPUT_PPTX = TOPIC_DIR / "skywork_exports" / f"{TOPIC}_skywork_v1.pptx"
OUTPUT_PDF = TOPIC_DIR / "skywork_exports" / f"{TOPIC}_skywork_v1.pdf"

FONT = "Malgun Gothic"
INK = RGBColor(25, 32, 29)
MUTED = RGBColor(86, 96, 91)
GREEN = RGBColor(15, 118, 110)
DARK_GREEN = RGBColor(31, 84, 61)
LIGHT_GREEN = RGBColor(230, 244, 241)
PALE_GREEN = RGBColor(244, 249, 247)
BLUE = RGBColor(37, 99, 235)
ORANGE = RGBColor(194, 65, 12)
RED = RGBColor(185, 28, 28)
GRAY = RGBColor(229, 234, 232)
WHITE = RGBColor(255, 255, 255)


def delete_template_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
    for slide_id in list(slide_id_list):
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def add_slide(prs: Presentation):
    layout_index = 6 if len(prs.slide_layouts) > 6 else 0
    return prs.slides.add_slide(prs.slide_layouts[layout_index])


def set_fill(shape, color: RGBColor, transparency: float = 0.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency
    shape.line.color.rgb = color


def add_rect(slide, x, y, w, h, color=WHITE, line=GRAY, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = line
    shp.line.width = Pt(0.75)
    return shp


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=14,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, *, size=12.4, color=INK, gap=2):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.space_after = Pt(gap)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def add_header(slide, title, subtitle=None, num=None):
    add_rect(slide, 0, 0, 13.333, 0.18, GREEN, GREEN)
    add_text(slide, 0.45, 0.34, 8.6, 0.34, title, size=20, bold=True, color=INK)
    if subtitle:
        add_text(slide, 0.47, 0.82, 8.95, 0.35, subtitle, size=10.5, bold=True, color=DARK_GREEN)
    if num:
        add_text(slide, 12.1, 0.34, 0.75, 0.28, f"{num:02d}", size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)


def add_footer(slide, text):
    add_text(slide, 0.45, 7.16, 11.95, 0.22, text, size=7.6, color=MUTED, valign=MSO_ANCHOR.MIDDLE)


def add_chip(slide, x, y, text, color=GREEN, width=1.15):
    shp = add_rect(slide, x, y, width, 0.32, LIGHT_GREEN, LIGHT_GREEN, radius=True)
    add_text(slide, x + 0.08, y + 0.06, width - 0.16, 0.18, text, size=8.2, bold=True, color=color, align=PP_ALIGN.CENTER)
    return shp


def add_card(slide, x, y, w, h, title, body, *, accent=GREEN, body_size=10.5):
    add_rect(slide, x, y, w, h, WHITE, GRAY, radius=False)
    add_rect(slide, x, y, 0.06, h, accent, accent)
    add_text(slide, x + 0.18, y + 0.17, w - 0.34, 0.28, title, size=11.2, bold=True, color=accent)
    if isinstance(body, list):
        add_bullets(slide, x + 0.18, y + 0.55, w - 0.34, h - 0.68, body, size=body_size, color=INK, gap=1)
    else:
        add_text(slide, x + 0.18, y + 0.55, w - 0.34, h - 0.68, body, size=body_size, color=INK)


def add_arrow(slide, x1, y1, x2, y2, color=GREEN):
    width = max(0.04, abs(x2 - x1))
    height = 0.035
    x = min(x1, x2)
    y = min(y1, y2) - height / 2
    return add_rect(slide, x, y, width, height, color, color)


def slide_title(prs):
    slide = add_slide(prs)
    add_rect(slide, 0, 0, 13.333, 7.5, WHITE, WHITE)
    add_rect(slide, 0, 0, 13.333, 0.22, GREEN, GREEN)
    add_text(slide, 0.65, 0.72, 3.6, 0.35, "AI TECH REVIEW | 2026-04-23", size=10.5, bold=True, color=GREEN)
    add_text(slide, 0.65, 1.45, 10.7, 0.95, "2026-04-17 AI Updates Weekly", size=33, bold=True, color=INK)
    add_text(
        slide,
        0.68,
        2.55,
        10.9,
        0.72,
        "모델 성능 경쟁에서 에이전트 운영 레이어 경쟁으로 이동하는 신호",
        size=19,
        bold=True,
        color=DARK_GREEN,
    )
    add_text(
        slide,
        0.7,
        3.62,
        8.4,
        1.18,
        "기반 자료: Lev Selector 영상, 동반 PPTX, 공식 Anthropic/OpenAI 문서, MCP docs, GitHub repos, 노동시장 tracker",
        size=13,
        color=MUTED,
    )
    for i, label in enumerate(["Model", "Harness", "Protocol", "Memory", "Governance"]):
        add_chip(slide, 0.7 + i * 1.45, 5.24, label, width=1.25)
    add_footer(slide, "Template reference: LGD_Template.pptx | Package: 2026-04-17_ai-updates-weekly")


def slide_signal_map(prs):
    slide = add_slide(prs)
    add_header(slide, "이번 업데이트의 증거 지형", "영상과 PPTX는 discovery layer, 결론은 공식 문서와 repo 검증 기준", 2)
    cols = [
        ("Confirmed", GREEN, ["Opus 4.7: 2026-04-16 공식 공개", "Agent SDK: tools/MCP/hooks/subagents", "MCP primitives: tools/resources/prompts", "Plugin marketplace docs", "Mythos Preview: gated cyber preview"]),
        ("Strong Secondary", BLUE, ["Anthropic $30B ARR signal", "OpenAI $25B 비교", "Broadcom/Google compute 규모", "일부 enterprise customer count"]),
        ("Weak / Unverified", RED, ["Mythos internal flywheel", "금융당국 emergency meeting", "전통 SW 즉시 소멸", "단일 노동시장 붕괴 서사"]),
        ("Interpretation", ORANGE, ["경쟁축은 agent operating layer", "프롬프트 -> 패키지/프로토콜", "개인 digital employee stack", "권한/감사가 adoption 병목"]),
    ]
    x = 0.55
    for title, color, items in cols:
        add_card(slide, x, 1.45, 2.95, 4.65, title, items, accent=color, body_size=9.6)
        x += 3.15
    add_footer(slide, "Sources: Anthropic Opus/SDK docs, MCP architecture, Anthropic GitHub repos, OpenAI enterprise note, source deck link harvest")


def slide_operating_layer(prs):
    slide = add_slide(prs)
    add_header(slide, "Agent Operating Layer의 구조", "더 좋은 답변보다 더 안정적인 실행면이 경쟁축", 3)
    add_card(slide, 0.65, 1.25, 3.75, 1.2, "Model", "Opus 4.7 / GPT-5.4 / Gemini / Mythos", accent=GREEN, body_size=10.2)
    add_card(slide, 4.8, 1.25, 3.75, 1.2, "Harness", "Claude Code / Agent SDK / desktop agents", accent=BLUE, body_size=10.2)
    add_card(slide, 8.95, 1.25, 3.75, 1.2, "Protocol", "MCP tools / resources / prompts", accent=ORANGE, body_size=10.2)
    add_card(slide, 0.65, 3.0, 3.75, 1.2, "Packaging", "plugins / skills / marketplaces", accent=GREEN, body_size=10.2)
    add_card(slide, 4.8, 3.0, 3.75, 1.2, "Memory", "Cognee / Obsidian / session summary", accent=BLUE, body_size=10.2)
    add_card(slide, 8.95, 3.0, 3.75, 1.2, "Execution", "CLI / browser / files / SaaS / deployment", accent=ORANGE, body_size=10.2)
    add_card(slide, 2.75, 4.75, 7.65, 0.95, "Governance", "permissions / audit / approval / sandbox / data boundaries", accent=RED, body_size=10.5)
    add_card(
        slide,
        0.65,
        6.25,
        11.15,
        0.82,
        "핵심 판단",
        "팀은 모델 leaderboard만 보지 말고, 도구 권한·메모리·감사·실패 복구·배포비용을 함께 평가해야 한다.",
        accent=ORANGE,
        body_size=9.6,
    )
    add_footer(slide, "Interpretation based on source pack + Anthropic/OpenAI/MCP docs")


def slide_anthropic_stack(prs):
    slide = add_slide(prs)
    add_header(slide, "Anthropic Stack: Opus 4.7 + Agent SDK + Claude Code", "모델 출시와 개발자 실행면이 같은 방향으로 움직임", 4)
    add_card(
        slide,
        0.55,
        1.35,
        3.9,
        4.5,
        "Claude Opus 4.7",
        ["2026-04-16 official announcement", "1M context window", "coding / AI agents / enterprise workflows", "API model id: claude-opus-4-7", "premium model: hard tasks first"],
        accent=GREEN,
        body_size=10.2,
    )
    add_card(
        slide,
        4.7,
        1.35,
        3.9,
        4.5,
        "Claude Agent SDK",
        ["Built-in tools: Read / Write / Edit / Bash", "Hooks, permissions, sessions", "Subagents for parallel/context isolation", "MCP integration", "External agents reuse Claude Code harness"],
        accent=BLUE,
        body_size=10.2,
    )
    add_card(
        slide,
        8.85,
        1.35,
        3.9,
        4.5,
        "Claude Code / Desktop",
        ["Releases around Apr 16-18 confirm fast iteration", "Developer-focused UI and routines signal", "Agent teams, browser preview, PR handling", "Execution surface, not only chat surface"],
        accent=ORANGE,
        body_size=10.2,
    )
    add_footer(slide, "Sources: anthropic.com/claude/opus, code.claude.com/docs/en/agent-sdk/overview, github.com/anthropics/claude-code/releases")


def slide_plugins_mcp(prs):
    slide = add_slide(prs)
    add_header(slide, "Plugins + MCP: 프롬프트가 패키지와 프로토콜로 변한다", "재사용 가능한 workflow는 instruction + tools + permissions + tests의 묶음", 5)
    add_card(
        slide,
        0.6,
        1.32,
        4.0,
        4.78,
        "Plugin / Skill Packaging",
        ["knowledge-work-plugins", "financial-services-plugins", ".claude-plugin/plugin.json", "marketplace.json", "repo/local cache 기반 배포", "팀별 workflow standardization"],
        accent=GREEN,
        body_size=10.2,
    )
    add_card(
        slide,
        4.85,
        1.32,
        3.6,
        4.78,
        "MCP Architecture",
        ["Host: Claude Code / Desktop / IDE", "Client: server별 dedicated connection", "Server: context provider", "Primitives: tools, resources, prompts", "Transport: stdio or streamable HTTP"],
        accent=BLUE,
        body_size=10.0,
    )
    add_card(
        slide,
        8.7,
        1.32,
        3.95,
        4.78,
        "왜 중요한가",
        ["integration code 반복을 줄임", "tool discovery와 call 방식 표준화", "agent가 SaaS/DB/files를 다룰 수 있음", "동시에 supply-chain surface도 확장"],
        accent=ORANGE,
        body_size=10.2,
    )
    add_footer(slide, "Sources: code.claude.com/docs/en/plugin-marketplaces, modelcontextprotocol.io/docs/learn/architecture")


def slide_mythos(prs):
    slide = add_slide(prs)
    add_header(slide, "Mythos Preview: Capability와 Control의 경계", "일반 공개가 아니라 보안 방어 중심 제한 preview로 읽어야 함", 6)
    add_text(slide, 0.7, 1.35, 5.35, 0.38, "공식 자료로 확인되는 내용", size=13, bold=True, color=DARK_GREEN)
    add_bullets(
        slide,
        0.8,
        1.9,
        5.45,
        2.15,
        [
            "containerized scaffold 안에서 source를 읽고 취약점 가설을 세움",
            "실행·debug·검증을 반복해 bug report와 PoC를 생성",
            "최종 triage와 coordinated disclosure 원칙을 강조",
            "general availability가 아니라 restricted preview",
        ],
        size=11.3,
    )
    add_text(slide, 6.85, 1.35, 5.5, 0.38, "슬라이드에서 낮춰야 할 표현", size=13, bold=True, color=RED)
    add_bullets(
        slide,
        6.95,
        1.9,
        5.45,
        2.15,
        [
            "internal flywheel / regulator emergency meeting는 official baseline 아님",
            "4x productivity 같은 내부 생산성 수치는 별도 확인 필요",
            "cyber 모델은 '더 강한 코딩 모델'보다 release governance 이슈",
        ],
        size=11.3,
    )
    add_rect(slide, 0.75, 4.65, 11.65, 1.05, PALE_GREEN, GRAY)
    add_text(slide, 1.0, 4.93, 11.1, 0.45, "Implication: 방어팀은 현재 공개 frontier model로 vulnerability finding practice를 시작하되, tool 권한·sandbox·audit·disclosure process를 먼저 설계해야 한다.", size=13, bold=True, color=INK)
    add_footer(slide, "Sources: red.anthropic.com/2026/mythos-preview, Claude Mythos Preview system card")


def slide_personal_stack(prs):
    slide = add_slide(prs)
    add_header(slide, "Personal Digital Employee Stack", "개인 자동화와 enterprise agent stack이 같은 재료를 사용하기 시작", 7)
    nodes = [
        ("Interface", "Claude Code\nOpenClaw\nHermes"),
        ("Context", "Obsidian\nGWS\nlocal files"),
        ("Tools", "MCP\nCLI\nbrowser"),
        ("Memory", "Cognee\ngraph\nsession"),
        ("Delivery", "Telegram\nSlack\nPR/docs"),
    ]
    x = 0.65
    for i, (label, body) in enumerate(nodes):
        add_rect(slide, x, 1.55, 2.15, 1.2, LIGHT_GREEN if i % 2 == 0 else WHITE, GRAY)
        add_text(slide, x + 0.15, 1.75, 1.75, 0.22, label, size=11.5, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.15, 2.1, 1.75, 0.46, body, size=9.5, color=INK, align=PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            add_arrow(slide, x + 2.17, 2.15, x + 2.45, 2.15, GREEN)
        x += 2.42
    add_card(
        slide,
        0.75,
        3.55,
        5.65,
        2.3,
        "Google Workspace CLI",
        ["Drive/Gmail/Calendar/Sheets/Docs/Chat", "structured JSON output", "agent skills and recipes", "OAuth setup is still the real adoption friction"],
        accent=BLUE,
        body_size=10.4,
    )
    add_card(
        slide,
        6.75,
        3.55,
        5.65,
        2.3,
        "Memory Layer",
        ["Cognee, Obsidian wiki/RAG, markdown graph", "vector-only retrieval is not enough for long-lived agents", "pinned facts + provenance + update policy needed"],
        accent=ORANGE,
        body_size=10.4,
    )
    add_footer(slide, "Sources: github.com/googleworkspace/cli, github.com/topoteretes/cognee, source deck")


def slide_workflow_deployment(prs):
    slide = add_slide(prs)
    add_header(slide, "Workflow UI와 Deployment Layer", "agent가 만든 결과를 실제 운영으로 넘기는 마지막 표면", 8)
    add_card(
        slide,
        0.6,
        1.3,
        3.8,
        4.65,
        "Workflow Builder",
        ["Langflow: runtime 포함", "React Flow: React 기반 canvas UI", "LiteGraph.js: vanilla JS / low overhead", "목적: human reviewable execution graph"],
        accent=GREEN,
        body_size=10.4,
    )
    add_card(
        slide,
        4.72,
        1.3,
        3.8,
        4.65,
        "Railway",
        ["GitHub repo / Docker image deploy", "build, networking, SSL, env config", "Hobby: $5 minimum usage", "deployment platform, not vibe-coding tool"],
        accent=BLUE,
        body_size=10.4,
    )
    add_card(
        slide,
        8.84,
        1.3,
        3.8,
        4.65,
        "Operating Question",
        ["누가 approve 하는가?", "실패 시 rollback은?", "로그와 비용은 어디에 남나?", "민감 데이터는 어디까지 나가나?", "model/backend swap 가능한가?"],
        accent=ORANGE,
        body_size=10.4,
    )
    add_footer(slide, "Sources: railway.com/pricing, langflow.org, reactflow.dev, github.com/jagenjo/litegraph.js")


def slide_market(prs):
    slide = add_slide(prs)
    add_header(slide, "Enterprise AI Market Signal", "수치보다 중요한 것은 양사 모두 operating layer로 이동한다는 점", 9)
    add_card(
        slide,
        0.65,
        1.4,
        5.8,
        3.75,
        "Anthropic Signal",
        ["$30B ARR claim: secondary market signal", "enterprise-heavy revenue narrative", "Opus / Code / SDK / plugins / Mythos가 한 방향", "official audited disclosure처럼 쓰면 안 됨"],
        accent=GREEN,
        body_size=10.8,
    )
    add_card(
        slide,
        6.85,
        1.4,
        5.8,
        3.75,
        "OpenAI Official Signal",
        ["2026-04-08 enterprise note", "enterprise = more than 40% of revenue", "consumer parity by end-2026 target", "Frontier + unified AI superapp + agents across systems"],
        accent=BLUE,
        body_size=10.8,
    )
    add_rect(slide, 0.75, 5.65, 11.8, 0.72, PALE_GREEN, GRAY)
    add_text(slide, 1.0, 5.87, 11.25, 0.25, "결론: Anthropic vs OpenAI의 실전 경쟁은 chat UX가 아니라 enterprise work surface, agent runtime, connectors, memory, governance에서 벌어진다.", size=12.2, bold=True, color=INK)
    add_footer(slide, "Sources: OpenAI next phase of enterprise AI, PYMNTS/Axios/Bloomberg-cited secondary reporting")


def slide_labor(prs):
    slide = add_slide(prs)
    add_header(slide, "Labor Market: 단일 서사로 읽으면 틀린다", "해고와 job posting 회복은 동시에 존재할 수 있음", 10)
    add_card(
        slide,
        0.65,
        1.35,
        3.75,
        4.7,
        "Layoffs remain elevated",
        ["Layoffs.fyi: 2026 tech layoffs continued", "TrueUp: impacted people and layoff events tracked", "AI Job Loss Tracker: AI material-factor methodology", "methodology마다 범위가 다름"],
        accent=RED,
        body_size=10.3,
    )
    add_card(
        slide,
        4.8,
        1.35,
        3.75,
        4.7,
        "Job postings can recover",
        ["Gizmodo/TrueUp: software job listings +30% YTD", "about 67,000 openings cited", "postings != actual hires", "ghost jobs / automated HR caveat"],
        accent=BLUE,
        body_size=10.3,
    )
    add_card(
        slide,
        8.95,
        1.35,
        3.5,
        4.7,
        "Practical reading",
        ["AI가 모든 개발자를 즉시 대체한다는 주장은 약함", "역할은 agent design/review/verification으로 이동", "교육 기준은 tool fluency + judgment + governance"],
        accent=GREEN,
        body_size=10.3,
    )
    add_footer(slide, "Sources: layoffs.fyi, trueup.io/layoffs, jobloss.ai, Gizmodo Apr 6 2026")


def slide_risks(prs):
    slide = add_slide(prs)
    add_header(slide, "Adoption Risks", "agent stack은 생산성만큼 blast radius도 키운다", 11)
    risks = [
        ("Vendor lock-in", "SDK + plugin + marketplace + memory가 한 벤더에 묶임"),
        ("Hidden cost", "1M context, long sessions, tool loops의 비용 예측 문제"),
        ("Permission sprawl", "files / browser / email / SaaS 권한이 확장"),
        ("Memory contamination", "outdated facts, private data, bad summaries 축적"),
        ("Tool supply chain", "lookalike tools, prompt injection, exfiltration"),
        ("Benchmark distraction", "leaderboard가 workflow success를 대체하지 못함"),
    ]
    for i, (title, body) in enumerate(risks):
        row = i // 2
        col = i % 2
        add_card(slide, 0.75 + col * 6.05, 1.25 + row * 1.63, 5.65, 1.15, title, body, accent=ORANGE if i % 2 else RED, body_size=9.8)
    add_footer(slide, "Risk framing based on report synthesis and MCP/tool/agent architecture constraints")


def slide_action_plan(prs):
    slide = add_slide(prs)
    add_header(slide, "30-90 Day Action Plan", "작은 workflow부터 권한과 검증 루프를 포함해 실험", 12)
    add_card(
        slide,
        0.65,
        1.3,
        3.85,
        4.9,
        "30 Days",
        ["반복 업무 10개 분류", "read-only / draft-only / approval / execute 권한 단계 정의", "Claude Code/Codex/Gemini/OpenClaw 중 하나로 작은 workflow 재현", "Obsidian/Drive/GitHub context 정리"],
        accent=GREEN,
        body_size=10.0,
    )
    add_card(
        slide,
        4.75,
        1.3,
        3.85,
        4.9,
        "60 Days",
        ["MCP/CLI 기반 실제 시스템 1개 연결", "session memory + run log 저장", "팀 공통 skills/plugins 3-5개 작성", "small service를 Railway/Vercel/Cloudflare/Render 중 하나에 deploy"],
        accent=BLUE,
        body_size=10.0,
    )
    add_card(
        slide,
        8.85,
        1.3,
        3.85,
        4.9,
        "90 Days",
        ["end-to-end agent workflow SOP화", "prompt injection / permissions / secrets review", "model/backend portability test", "cycle time / rework / defect escape 지표화"],
        accent=ORANGE,
        body_size=10.0,
    )
    add_footer(slide, "Action plan adapted from deep research report")


def slide_pilot_architecture(prs):
    slide = add_slide(prs)
    add_header(slide, "Pilot Architecture: Agent Workflow의 최소 운영 골격", "실험이라도 governance를 나중에 붙이지 말고 처음부터 넣는다", 13)
    boxes = [
        ("Intake", "topic / ticket / doc"),
        ("Context", "source pack\nrepo / drive / obsidian"),
        ("Agent", "model + SDK\npermissions"),
        ("Tools", "MCP / CLI / browser\nread + write gates"),
        ("Verify", "tests / lint / human review"),
        ("Publish", "report / PR / deck / deploy"),
    ]
    x = 0.55
    for i, (t, b) in enumerate(boxes):
        add_rect(slide, x, 1.7, 1.75, 1.15, LIGHT_GREEN if i in {0, 2, 4} else WHITE, GRAY)
        add_text(slide, x + 0.12, 1.92, 1.5, 0.2, t, size=10.5, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.12, 2.22, 1.5, 0.32, b, size=8.8, color=INK, align=PP_ALIGN.CENTER)
        if i < len(boxes) - 1:
            add_arrow(slide, x + 1.78, 2.27, x + 2.08, 2.27)
        x += 2.05
    add_card(
        slide,
        0.8,
        4.05,
        3.55,
        1.55,
        "Must log",
        ["source IDs", "tool calls", "permissions used", "cost/session time", "human approvals"],
        accent=GREEN,
        body_size=9.6,
    )
    add_card(
        slide,
        4.75,
        4.05,
        3.55,
        1.55,
        "Must block",
        ["secret export", "unreviewed write actions", "unknown MCP tools", "network calls from unsafe context"],
        accent=RED,
        body_size=9.6,
    )
    add_card(
        slide,
        8.7,
        4.05,
        3.55,
        1.55,
        "Must measure",
        ["cycle time", "review time", "defects caught", "rollback count", "user trust"],
        accent=BLUE,
        body_size=9.6,
    )
    add_footer(slide, "Recommended pilot pattern for engineering and product teams")


def slide_final(prs):
    slide = add_slide(prs)
    add_header(slide, "Final Synthesis", "좋은 모델을 사는 것보다 좋은 실행면을 설계하는 팀이 이긴다", 14)
    add_text(
        slide,
        0.8,
        1.42,
        11.6,
        0.85,
        "AI 경쟁은 모델 지능 단독 경쟁에서, 에이전트를 실제 업무 맥락 안에서 안전하게 실행시키는 운영 레이어 경쟁으로 이동하고 있다.",
        size=22,
        bold=True,
        color=DARK_GREEN,
        align=PP_ALIGN.CENTER,
    )
    add_card(
        slide,
        1.1,
        3.05,
        3.45,
        2.1,
        "Engineering",
        ["MCP/CLI/tool permission 설계", "verification loop와 rollback", "memory provenance 관리"],
        accent=GREEN,
        body_size=10.2,
    )
    add_card(
        slide,
        4.95,
        3.05,
        3.45,
        2.1,
        "Product",
        ["agent workflow를 feature가 아니라 operating flow로 정의", "human approval UX 설계", "trust and audit as product surface"],
        accent=BLUE,
        body_size=10.2,
    )
    add_card(
        slide,
        8.8,
        3.05,
        3.45,
        2.1,
        "Strategy",
        ["vendor portability 유지", "enterprise AI budget은 workflow ROI로 평가", "labor role redesign 준비"],
        accent=ORANGE,
        body_size=10.2,
    )
    add_footer(slide, "See reports/2026-04-17_ai-updates-weekly_deepresearch.md for full evidence and references")


def export_pdf(pptx_path: Path, pdf_path: Path) -> bool:
    try:
        import pythoncom
        import win32com.client  # type: ignore
    except Exception:
        return False

    app = None
    presentation = None
    try:
        pythoncom.CoInitialize()
        app = win32com.client.Dispatch("PowerPoint.Application")
        app.Visible = 1
        presentation = app.Presentations.Open(str(pptx_path), WithWindow=False)
        presentation.SaveAs(str(pdf_path), 32)
        return True
    except Exception as exc:
        print(f"pdf_export_error={exc}")
        return False
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def main() -> int:
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    for builder in [
        slide_title,
        slide_signal_map,
        slide_operating_layer,
        slide_anthropic_stack,
        slide_plugins_mcp,
        slide_mythos,
        slide_personal_stack,
        slide_workflow_deployment,
        slide_market,
        slide_labor,
        slide_risks,
        slide_action_plan,
        slide_pilot_architecture,
        slide_final,
    ]:
        builder(prs)

    prs.save(str(OUTPUT_PPTX))
    pdf_ok = export_pdf(OUTPUT_PPTX, OUTPUT_PDF)
    print(f"saved_pptx={OUTPUT_PPTX}")
    print(f"saved_pdf={OUTPUT_PDF if pdf_ok else 'NOT_CREATED'}")
    return 0 if pdf_ok else 1


if __name__ == "__main__":
    raise SystemExit(main())
