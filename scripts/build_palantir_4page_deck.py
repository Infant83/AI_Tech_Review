from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\angpa\myProjects\Daily_Work\AI_Tech_Review")
TOPIC_DIR = ROOT / "2026-04-15_palantir"
TEMPLATE = TOPIC_DIR / "sources" / "ppt_template" / "Template_4pages_new.pptx"
OUTPUT_PPTX = TOPIC_DIR / "skywork_exports" / "2026-04-16_palantir_4page_template_new_local.pptx"
OUTPUT_PDF = TOPIC_DIR / "skywork_exports" / "2026-04-16_palantir_4page_template_new_local.pdf"

FONT = "Malgun Gothic"
COLOR_TEXT = RGBColor(32, 33, 36)
COLOR_SUB = RGBColor(31, 84, 61)
COLOR_FOOT = RGBColor(90, 96, 104)
COLOR_WHITE = RGBColor(255, 255, 255)


def shp(slide, idx: int):
    return slide.shapes[idx - 1]


def set_text(
    shape,
    text: str,
    *,
    size: float = 13,
    bold: bool = False,
    color: RGBColor = COLOR_TEXT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    paragraphs = text.split("\n")
    for i, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: float = 11,
    bold: bool = False,
    color: RGBColor = COLOR_TEXT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    set_text(box, text, size=size, bold=bold, color=color, align=align, valign=valign)
    return box


def set_circle_label(shape, text: str):
    set_text(
        shape,
        text,
        size=22,
        bold=True,
        color=COLOR_WHITE,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_footer(slide, text: str):
    add_textbox(
        slide,
        0.42,
        7.08,
        12.35,
        0.24,
        text,
        size=8.5,
        color=COLOR_FOOT,
        valign=MSO_ANCHOR.MIDDLE,
    )


def fill_slide_1(slide):
    set_text(shp(slide, 1), "보안·데이터주권·온톨로지 구축", size=24, bold=True)
    set_text(
        shp(slide, 2),
        "고위험 제조데이터에서도 핵심은 AI 사용 여부가 아니라 어디에 두고 어떤 권한·행위통제로 운영하느냐다",
        size=14,
        bold=True,
        color=COLOR_SUB,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_text(
        shp(slide, 3),
        "\n".join(
            [
                "권장 아키텍처",
                "• 원천 데이터: ERP / MES / QMS / EAM / Sensor / Recipe / Defect / Safety",
                "• 고객 통제 Ontology: lot / panel / tool / defect / recipe / alarm / action",
                "• 실행 영역: ① On-prem·air-gapped ② Private cloud ③ 승인 저위험 데이터만 외부 LLM",
                "• 원천 데이터·권한·export policy는 고객 유지, Palantir는 workflow·governance를 공동 설계",
            ]
        ),
        size=13.2,
    )
    set_text(
        shp(slide, 4),
        "\n".join(
            [
                "보안 관점 정리",
                "• 장점: raw 제조데이터 public SaaS 미반출, lineage·checkpoint·audit, object/action 권한, private path 설계",
                "• 남는 과제: on-prem = zero risk 아님, support 계정·patch·logging·export path 관리, ontology 공수",
                "• 대안 비교: Palantir = ontology+workflow+governance 통합 / Azure·Vertex·Bedrock = private AI 가능하나 거버넌스 조합 필요 / In-house GraphRAG = 통제 강하나 구축·유지부담 큼",
                "• LGD 포인트: 'cloud 금지'보다 '어떤 데이터와 action을 밖에 둘 수 없는가'를 먼저 정의",
            ]
        ),
        size=12.2,
    )
    add_footer(
        slide,
        "출처: Palantir Apollo · Palantir Privacy and Governance Whitepaper · Azure data privacy / private networking · Vertex AI VPC-SC · law.go.kr",
    )


def fill_slide_2(slide):
    set_text(shp(slide, 1), "경쟁사·시장 신호와 Samsung 흐름", size=24, bold=True)
    set_text(
        shp(slide, 2),
        "2025-2026 제조 AI 경쟁은 chatbot이 아니라 ontology·workflow·digital twin·operations intelligence로 이동 중",
        size=14,
        bold=True,
        color=COLOR_SUB,
        valign=MSO_ANCHOR.MIDDLE,
    )

    set_text(shp(slide, 4), "추진 배경", size=12.5, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 6), "25년까지 공개 사실", size=12.5, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 32), "26년 이후 신호", size=12.5, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    left_boxes = {
        9: "챗봇보다 공정·품질·물류 등 운영형 AI 수요 확대",
        8: "문제는 모델보다 ontology·workflow·승인통제 설계",
        12: "품질·수율·설비·SCM처럼 cross-functional 문제로 이동",
        14: "채용 신호도 KG·GraphRAG·NL2SQL로 이동",
        33: "핵심은 customer AI가 아니라 operations intelligence",
    }
    mid_boxes = {
        16: "2024-08 DX가 MS·Google·Palantir 3사 PoC",
        18: "2024-09 대고객 rollout은 Microsoft 선정",
        22: "2025 AI Forum: 제조데이터 체계화·고품질화 강조",
        20: "삼성 = Palantir 성공사례로 볼 공개 production proof는 부족",
    }
    right_boxes = {
        24: "2026-02 DS AI Center KG JD: Palantir ontology·Neo4j·GraphRAG",
        26: "2026-03 GTC: Agentic AI + digital twin 공개축은 NVIDIA·Synopsys",
        30: "2026 LG CNS partnership, HD Hyundai group-wide expansion",
        28: "2nm yield 개선 신호는 존재하나 Palantir attribution은 미확인",
    }
    for idx, text in left_boxes.items():
        set_text(shp(slide, idx), text, size=10.4, valign=MSO_ANCHOR.MIDDLE)
    for idx, text in mid_boxes.items():
        set_text(shp(slide, idx), text, size=10.4, valign=MSO_ANCHOR.MIDDLE)
    for idx, text in right_boxes.items():
        set_text(shp(slide, idx), text, size=10.1, valign=MSO_ANCHOR.MIDDLE)

    add_footer(
        slide,
        "공개 fact와 signal 분리: ETNews 2024-08-05 / 2024-09-02 · Samsung AI Forum 2025 · Samsung GTC 2026 · Hankyung 2026-03-31 · LG CNS · Reuters/Yahoo HD Hyundai",
    )


def fill_slide_3(slide):
    set_text(shp(slide, 1), "계열사 readiness 및 비용 검토", size=24, bold=True)
    set_text(
        shp(slide, 2),
        "계열사 readiness는 동일하지 않다. 현 단계에서는 stage·문제영역·비용구조를 함께 보는 working outline이 필요하다",
        size=13.5,
        bold=True,
        color=COLOR_SUB,
        valign=MSO_ANCHOR.MIDDLE,
    )

    set_text(shp(slide, 11), "단계별 도입 전략", size=12.5, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 12), "", size=10)
    set_text(shp(slide, 14), "계열사 적용 현황", size=12.5, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    set_text(shp(slide, 15), "설명회 / 문제 정의", size=11.5, bold=True, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 16), "Pilot 적용 / 타당성", size=11.5, bold=True, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 17), "본사업 / 운영 적용", size=11.5, bold=True, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 18), "비용 / 검토 메모", size=11.5, bold=True, valign=MSO_ANCHOR.MIDDLE)

    set_text(shp(slide, 9), "대상 문제, 데이터 경계, 금지 데이터, owner 정의", size=12.2, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 7), "8~12주 KPI proof / ontology 최소모델 / 사용자 승인 flow", size=12.0, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 5), "라인·조직 확장 / workflow action 연결 / support 체계", size=12.0, valign=MSO_ANCHOR.MIDDLE)
    set_text(
        shp(slide, 3),
        "비용은 bootcamp + ontology + integration + expansion 구조\nstatus 일부는 내부회의 기반 추정이며 추가 검증 필요",
        size=11.8,
        valign=MSO_ANCHOR.MIDDLE,
    )

    set_text(shp(slide, 20), "LGES", size=12, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 22), "LGE", size=12, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 24), "LG Innotek", size=12, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 26), "LGD", size=12, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 36), "LG CNS", size=12, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    set_text(shp(slide, 52), "5개 PoC 언급", size=10.3, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 54), "quality 중심", size=10.3, valign=MSO_ANCHOR.MIDDLE)

    set_text(shp(slide, 51), "초기 검토 추정", size=10.3, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 53), "owner·stage 미확인", size=10.3, valign=MSO_ANCHOR.MIDDLE)

    set_text(shp(slide, 48), "PoC / 문제선정", size=10.3, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 47), "미도입 추정", size=10.3, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 44), "POC 8억 언급", size=10.0, valign=MSO_ANCHOR.MIDDLE)

    set_text(shp(slide, 49), "공식 partnership", size=10.1, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 50), "late-2025 deployment", size=9.9, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 46), "quality full-scale signal", size=9.8, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 45), "project expansion형", size=10.0, valign=MSO_ANCHOR.MIDDLE)

    add_textbox(slide, 8.24, 2.20, 1.05, 0.52, "검토 전 /\nPoC-first", size=10.4, valign=MSO_ANCHOR.MIDDLE)

    add_footer(
        slide,
        "해석 주의: LGES·LGE·LG Innotek 일부 status와 비용 메모는 내부회의 기반 추정. 공개 fact는 LG CNS partnership 중심이며, 이후 코멘트 추가를 전제로 한 working outline.",
    )


def fill_slide_4(slide):
    set_text(shp(slide, 1), "자사 도입방안 (LGD)", size=24, bold=True)
    set_text(
        shp(slide, 2),
        "도입 판단의 핵심은 어디에서 가장 빨리, 가장 안전하게, 가장 설명 가능한 운영 성과를 증명할 수 있는가다",
        size=13.3,
        bold=True,
        color=COLOR_SUB,
        valign=MSO_ANCHOR.MIDDLE,
    )

    set_text(shp(slide, 18), "검증 Criteria", size=12.5, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 38), "영역", size=10.8, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 40), "효과성", size=10.8, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 42), "검증", size=10.8, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 44), "비용부담", size=10.4, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 46), "확대성", size=10.8, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    set_text(shp(slide, 21), "품질/Q-Cost", size=10.8, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 22), "공정 이상", size=10.8, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 27), "안전", size=10.8, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 28), "설비 보전", size=10.4, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 32), "SCM/병목", size=10.8, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 35), "R&D/지식", size=10.8, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    ratings = {
        48: "상",
        50: "상",
        52: "중",
        54: "상",
        64: "상",
        66: "중상",
        68: "중",
        70: "상",
        56: "중상",
        58: "중상",
        60: "하",
        62: "중",
        72: "중상",
        74: "중",
        76: "중",
        78: "중상",
        88: "중상",
        90: "중",
        92: "중상",
        94: "상",
        80: "중",
        82: "중하",
        84: "중",
        86: "중상",
    }
    for idx, value in ratings.items():
        set_text(shp(slide, idx), value, size=10.5, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    set_text(shp(slide, 6), "우선 PoC 1", size=12.5, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 7), "품질 / Q-Cost", size=14, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 9), "우선 PoC 2", size=12.5, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_text(shp(slide, 10), "공정 이상 +\n원인분석", size=13.2, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_circle_label(shp(slide, 12), "1")
    set_circle_label(shp(slide, 13), "2")
    set_text(
        shp(slide, 14),
        "재무효과를 가장 설명하기 쉽고\nLG quality reference·현업 KPI와 연결 가능",
        size=10.8,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_text(
        shp(slide, 15),
        "lot·recipe·defect ontology에 적합하고\n확대성 높은 제조 workflow형 문제",
        size=10.8,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_text(shp(slide, 11), "3순위: 안전 leading indicator", size=10.8, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(
        slide,
        8.80,
        6.55,
        3.15,
        0.34,
        "체크포인트: 금지 데이터 범위 · 8~12주 KPI · ontology owner · human approval action",
        size=8.9,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        color=COLOR_FOOT,
    )
    add_footer(
        slide,
        "판단 기준: 효과성 / 검증 가능성 / 비용부담 / 확대성. 전사 rollout 논의보다 PoC 정의와 stop-go 기준 합의가 우선.",
    )


def try_export_pdf(pptx_path: Path, pdf_path: Path) -> bool:
    try:
        import pythoncom
        import win32com.client  # type: ignore
    except Exception:
        return False

    try:
        pythoncom.CoInitialize()
        app = win32com.client.Dispatch("PowerPoint.Application")
        app.Visible = 1
        presentation = app.Presentations.Open(str(pptx_path), WithWindow=False)
        presentation.SaveAs(str(pdf_path), 32)
        presentation.Close()
        app.Quit()
        return True
    except Exception:
        return False
    finally:
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def main():
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(TEMPLATE))

    fill_slide_1(prs.slides[0])
    fill_slide_2(prs.slides[1])
    fill_slide_3(prs.slides[2])
    fill_slide_4(prs.slides[3])

    prs.save(str(OUTPUT_PPTX))
    pdf_ok = try_export_pdf(OUTPUT_PPTX, OUTPUT_PDF)
    print(f"saved_pptx={OUTPUT_PPTX}")
    print(f"saved_pdf={OUTPUT_PDF if pdf_ok else 'NOT_CREATED'}")


if __name__ == "__main__":
    main()
