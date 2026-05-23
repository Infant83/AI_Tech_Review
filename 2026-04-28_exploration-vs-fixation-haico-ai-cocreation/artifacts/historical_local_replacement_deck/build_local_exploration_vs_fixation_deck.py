from __future__ import annotations

from pathlib import Path
import textwrap

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas


ROOT = Path(r"C:\Users\angpa\myProjects\Daily_Work\AI_Tech_Review")
TOPIC = ROOT / "2026-04-28_exploration-vs-fixation-haico-ai-cocreation"
TEMPLATE = Path(r"C:\Users\angpa\.codex\skills\skywork-ppt-workflow\assets\LGD_Template.pptx")
OUT_PPTX = TOPIC / "skywork_exports" / "2026-04-28_exploration-vs-fixation_skywork_local_v1.pptx"
OUT_PDF = TOPIC / "skywork_exports" / "2026-04-28_exploration-vs-fixation_skywork_local_v1.pdf"

FONT = "Malgun Gothic"
INK = RGBColor(29, 35, 33)
MUTED = RGBColor(86, 96, 91)
GREEN = RGBColor(15, 118, 110)
DARK_GREEN = RGBColor(31, 84, 61)
LIGHT_GREEN = RGBColor(229, 244, 241)
PALE = RGBColor(247, 250, 249)
ORANGE = RGBColor(194, 65, 12)
BLUE = RGBColor(37, 99, 235)
RED = RGBColor(185, 28, 28)
GRAY = RGBColor(219, 226, 223)
WHITE = RGBColor(255, 255, 255)

SLIDES = [
    {
        "title": "AI 협업은 첫 산출물을 늦출수록 넓어진다",
        "subtitle": "Exploration vs. Fixation 논문 기반 실무 리뷰",
        "body": [
            "Wen et al., arXiv:2512.18388v2, 2026-04-06",
            "리뷰 초점: 논문 인사이트, AI 사용자의 작업 순서, 조직 단위 다양성 관리",
        ],
        "footer": "Sources: Wen et al. 2026; AI Matters 2026-04-20",
        "layout": "title",
    },
    {
        "title": "읽어야 하는 이유",
        "subtitle": "문제는 모델 성능보다 사용 흐름에 가깝다",
        "cards": [
            ("빠른 결과", "챗봇은 요청 직후 완성품처럼 보이는 산출물을 만든다."),
            ("숨은 고정점", "첫 산출물은 이후 판단과 수정의 기준점으로 작동한다."),
            ("좁은 수정", "사용자는 새 방향을 찾기보다 기존 결과를 조금씩 고친다."),
            ("실무 위험", "보고서, 슬라이드, 코드, 디자인 모두 초기 방향에 묶일 수 있다."),
        ],
        "footer": "Primary source: arXiv:2512.18388v2, introduction and problem analysis",
    },
    {
        "title": "논문이 잡은 세 가지 병목",
        "subtitle": "창의 작업에서 바로 실행하는 인터페이스가 만드는 손실",
        "cards": [
            ("Premature convergence", "탐색 전에 하나의 방향으로 빨리 수렴한다."),
            ("Design fixation", "초기 예시나 결과물의 특징에 묶여 대안 폭이 줄어든다."),
            ("Gulf of envisioning", "원하는 느낌은 있지만 모델이 실행할 지시로 바꾸기 어렵다."),
        ],
        "flow": ["Prompt", "First artifact", "Local edits", "Good enough"],
        "footer": "Sources: Wen et al. 2026; Jansson and Smith 1991 cited in paper",
    },
    {
        "title": "HAICo의 설계",
        "subtitle": "탐색판과 실행판을 분리하고, 수정 의도를 실행 전에 보이게 한다",
        "cards": [
            ("Divergent mode", "이미지 생성 전에 idea grid를 만든다. 원격 개념, 역사, 신화, 인터넷 문화 같은 영역을 끌어온다."),
            ("Convergent mode", "선택한 방향을 semantic parameter와 option으로 분해한다. 사용자는 실행 전 해석을 본다."),
            ("Non-linear workflow", "탭과 라이브러리로 여러 갈래를 보존한다. 사용자는 돌아가고 비교하고 다시 조합한다."),
        ],
        "footer": "Primary source: arXiv:2512.18388v2, system design",
    },
    {
        "title": "실험 결과",
        "subtitle": "HAICo의 강점은 novelty와 diversity에 집중됐다. Fluency와 usefulness는 유의차가 없었다",
        "table": [
            ("Study", "within-subjects poster task", "N = 24"),
            ("CSI", "HAICo higher on all 5 dimensions", "all p < 0.002"),
            ("UMUX-Lite", "81.25 vs 64.24", "p < 0.001"),
            ("Novelty", "3.22 vs 2.41", "p < 0.001"),
            ("Diversity", "0.48 vs 0.36", "p = 0.001"),
            ("Fluency / usefulness", "no significant difference", "caution"),
        ],
        "footer": "Primary source: arXiv:2512.18388v2, Fig. 5 and section 6.1",
    },
    {
        "title": "사용자는 무엇을 배우는가",
        "subtitle": "도구 조작 학습에서 과제와 작업 순서 학습으로 이동한다",
        "cards": [
            ("ChatGPT condition", "System behaviors, prompting strategies가 많이 보고됐다. 도구를 어떻게 다뤄야 하는지에 주의가 간다."),
            ("HAICo condition", "Task-specific knowledge, new directions, workflow learning이 더 많이 보고됐다."),
            ("Measured signal", "Self-reported learning: HAICo 5.29, ChatGPT 3.12, p < 0.001."),
        ],
        "footer": "Primary source: arXiv:2512.18388v2, Fig. 9 and section 7.4",
    },
    {
        "title": "외부 근거가 보강하는 지점",
        "subtitle": "개인 생산성과 집단 다양성을 따로 관리해야 한다",
        "cards": [
            ("Design fixation with GenAI", "Wadinambiarachchi et al. 2024: AI 이미지 지원은 초기 예시 고착, 아이디어 수 감소, 다양성/독창성 저하와 연결됐다."),
            ("Individual gain, collective loss", "Doshi and Hauser 2024: GenAI 아이디어는 개인 글쓰기 평가를 높였지만 이야기 간 유사성도 높였다."),
            ("Homogenized ideation", "Anderson et al. 2024: ChatGPT 사용자는 상세 아이디어를 더 만들었지만 사용자 간 아이디어가 덜 구별됐다."),
        ],
        "footer": "Sources: arXiv:2403.11164; arXiv:2312.00506; arXiv:2402.01536",
    },
    {
        "title": "AI 사용 워크플로",
        "subtitle": "최종 산출물 전에 탐색과 선택을 명시한다",
        "flow": [
            "Problem frame",
            "Divergent cards",
            "Selection criteria",
            "Semantic parameters",
            "Artifact",
            "Audit",
        ],
        "bullets": [
            "첫 요청은 방향 카드로 시작한다.",
            "선택 기준을 점수와 위험으로 기록한다.",
            "수정 의도는 실행 전에 매개변수로 확인한다.",
            "버린 갈래와 이유를 보관한다.",
        ],
        "footer": "Derived workflow from Wen et al. 2026 and supporting studies",
    },
    {
        "title": "바로 쓸 수 있는 요청 패턴",
        "subtitle": "프롬프트보다 흐름이 중요하다",
        "cards": [
            ("Idea card pass", "서로 멀리 떨어진 아이디어 카드 9개를 title, source domain, concept, fit, risk로 작성."),
            ("Selection matrix", "novelty, usefulness, feasibility, audience fit, risk로 평가하고 선택 이유를 기록."),
            ("Parameter pass", "수정 의도를 semantic parameter, option, expected effect, side effect로 분해."),
            ("Team diversity", "팀원별 persona, source domain, 반대 가설을 나눠 중복을 줄임."),
        ],
        "footer": "Prompt patterns adapted from the report workflow",
    },
    {
        "title": "업무별 적용",
        "subtitle": "이미지 생성 연구지만 원리는 다른 작업에도 옮겨볼 수 있다",
        "table": [
            ("리서치", "요약 요청", "claim map, 반례, 후속 연구"),
            ("슬라이드", "10장 생성", "narrative 후보와 evidence map"),
            ("제품 기획", "기능 아이디어", "사용자 문제와 실패 시나리오 카드"),
            ("코딩", "바로 구현", "아키텍처 대안과 rollback path"),
            ("회의", "회의록 정리", "결정, 가정, 미결정 분리"),
        ],
        "footer": "Caution: cross-domain application remains an extrapolation from an image-generation study",
    },
    {
        "title": "제한과 남은 질문",
        "subtitle": "강한 결과지만 적용 범위는 좁게 읽어야 한다",
        "cards": [
            ("Study limits", "N=24, CS/IT 배경 치우침, 포스터 이미지 과제, 단회성 세션."),
            ("Learning limits", "학습 효과는 자기보고 중심이다. 장기 유지 효과는 아직 확인되지 않았다."),
            ("Product limits", "HAICo는 연구 시스템이다. 공개 도구보다 설계 원리를 가져오는 편이 현실적이다."),
            ("Adoption risk", "스캐폴딩은 목표가 명확한 작업에서 속도와 주도감을 낮출 수 있다."),
        ],
        "footer": "Primary source: arXiv:2512.18388v2, limitations and conclusion",
    },
    {
        "title": "운영 원칙",
        "subtitle": "AI를 많이 쓰는 팀일수록 결과물 다양성을 설계해야 한다",
        "cards": [
            ("1. Alternatives before artifacts", "AI가 완성품을 만들기 전에 서로 먼 대안을 먼저 보여주게 한다."),
            ("2. Interpretation before execution", "모호한 수정 요청은 실행 전에 매개변수와 옵션으로 확인한다."),
            ("3. Diversity at team level", "같은 모델, 같은 프롬프트, 같은 persona가 반복되지 않게 배정한다."),
            ("4. Archive rejected branches", "버린 대안과 이유를 남겨 다음 탐색의 출발점으로 쓴다."),
        ],
        "footer": "Operating rule for AI_Tech_Review and practical AI workflows",
    },
]


def delete_template_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001
    for slide_id in list(slide_id_list):
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def add_slide(prs: Presentation):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    return prs.slides.add_slide(layout)


def add_rect(slide, x, y, w, h, fill=WHITE, line=GRAY):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.65)
    return shp


def add_text(slide, x, y, w, h, text, size=13, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(0)
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return box


def add_header(slide, idx, title, subtitle):
    add_rect(slide, 0, 0, 13.333, 0.17, GREEN, GREEN)
    add_text(slide, 0.42, 0.32, 9.1, 0.35, title, 20, True, INK)
    add_text(slide, 0.44, 0.79, 9.8, 0.32, subtitle, 10.5, True, DARK_GREEN)
    add_text(slide, 12.18, 0.32, 0.72, 0.24, f"{idx:02d}", 10.5, True, GREEN, PP_ALIGN.RIGHT)


def add_footer(slide, text):
    add_text(slide, 0.45, 7.18, 12.2, 0.18, text, 7.3, False, MUTED)


def add_card(slide, x, y, w, h, title, body, accent=GREEN):
    add_rect(slide, x, y, w, h, WHITE, GRAY)
    add_rect(slide, x, y, 0.055, h, accent, accent)
    add_text(slide, x + 0.16, y + 0.14, w - 0.28, 0.25, title, 11.2, True, accent)
    add_text(slide, x + 0.16, y + 0.48, w - 0.28, h - 0.55, body, 10.2, False, INK)


def add_flow(slide, items, y=4.78):
    n = len(items)
    w = 11.8 / n
    x0 = 0.72
    for i, item in enumerate(items):
        x = x0 + i * w
        add_rect(slide, x, y, w - 0.13, 0.55, LIGHT_GREEN if i % 2 == 0 else PALE, GREEN)
        add_text(slide, x + 0.08, y + 0.17, w - 0.3, 0.2, item, 9.3, True, DARK_GREEN, PP_ALIGN.CENTER)
        if i < n - 1:
            add_rect(slide, x + w - 0.18, y + 0.26, 0.2, 0.035, GREEN, GREEN)


def add_bullets(slide, x, y, w, h, bullets, size=11.5):
    lines = "\n".join(f"- {b}" for b in bullets)
    add_text(slide, x, y, w, h, lines, size, False, INK)


def add_table(slide, rows, x=0.65, y=1.52, w=12.0, h=4.75):
    col_w = [0.25 * w, 0.49 * w, 0.26 * w]
    row_h = h / len(rows)
    for r, row in enumerate(rows):
        fill = PALE if r % 2 else WHITE
        cy = y + r * row_h
        cx = x
        for c, value in enumerate(row):
            add_rect(slide, cx, cy, col_w[c], row_h, fill, GRAY)
            add_text(slide, cx + 0.08, cy + 0.13, col_w[c] - 0.16, row_h - 0.14, value, 9.5 if r else 9.8, c == 0, GREEN if c == 0 else INK)
            cx += col_w[c]


def render_pptx():
    prs = Presentation(str(TEMPLATE)) if TEMPLATE.exists() else Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    delete_template_slides(prs)

    for idx, spec in enumerate(SLIDES, 1):
        slide = add_slide(prs)
        add_rect(slide, 0, 0, 13.333, 7.5, WHITE, WHITE)
        if spec.get("layout") == "title":
            add_rect(slide, 0, 0, 13.333, 0.22, GREEN, GREEN)
            add_text(slide, 0.7, 0.72, 3.2, 0.28, "AI TECH REVIEW | 2026-04-28", 10, True, GREEN)
            add_text(slide, 0.68, 1.5, 10.8, 0.9, spec["title"], 31, True, INK)
            add_text(slide, 0.72, 2.6, 9.8, 0.46, spec["subtitle"], 18, True, DARK_GREEN)
            add_bullets(slide, 0.76, 3.45, 9.7, 0.85, spec["body"], 12.3)
            add_rect(slide, 0.72, 5.0, 11.8, 0.78, PALE, GRAY)
            add_text(slide, 0.95, 5.19, 11.1, 0.3, "작업 규칙: 먼저 탐색하고, 선택 기준을 세운 뒤 실행한다.", 15, True, GREEN)
            add_footer(slide, spec["footer"])
            continue

        add_header(slide, idx, spec["title"], spec.get("subtitle", ""))
        if "cards" in spec:
            cards = spec["cards"]
            cols = 2 if len(cards) in (4, 3) else 3
            card_w = 5.9 if cols == 2 else 3.8
            start_x = 0.68
            start_y = 1.45
            for i, (title, body) in enumerate(cards):
                row = i // cols
                col = i % cols
                x = start_x + col * (card_w + 0.35)
                y = start_y + row * 1.62
                h = 1.32 if len(cards) > 3 else 1.45
                add_card(slide, x, y, card_w, h, title, body, [GREEN, BLUE, ORANGE, RED][i % 4])
        if "table" in spec:
            add_table(slide, spec["table"])
        if "flow" in spec:
            add_flow(slide, spec["flow"], 5.15 if "bullets" in spec else 5.55)
        if "bullets" in spec:
            add_bullets(slide, 0.9, 3.95, 10.6, 0.8, spec["bullets"], 11.2)
        add_footer(slide, spec["footer"])

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)


def register_pdf_font():
    font_dir = Path(r"C:\Windows\Fonts")
    regular = font_dir / "malgun.ttf"
    bold = font_dir / "malgunbd.ttf"
    if regular.exists():
        pdfmetrics.registerFont(TTFont("Malgun", str(regular)))
    if bold.exists():
        pdfmetrics.registerFont(TTFont("MalgunBold", str(bold)))


def rgb_to_hex(rgb):
    return colors.HexColor(f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")


def draw_text(c, x, y, w, text, size=11, bold=False, color=INK, leading=None):
    font = "MalgunBold" if bold else "Malgun"
    c.setFont(font, size)
    c.setFillColor(rgb_to_hex(color))
    leading = leading or size * 1.35
    max_chars = max(8, int(w / (size * 0.42)))
    lines = []
    for raw in str(text).split("\n"):
        if not raw:
            lines.append("")
        else:
            lines.extend(textwrap.wrap(raw, width=max_chars, break_long_words=False, replace_whitespace=False))
    yy = y
    for line in lines:
        c.drawString(x, yy, line)
        yy -= leading


def draw_rect(c, x, y, w, h, fill=WHITE, stroke=GRAY):
    c.setFillColor(rgb_to_hex(fill))
    c.setStrokeColor(rgb_to_hex(stroke))
    c.rect(x, y, w, h, fill=1, stroke=1)


def render_pdf():
    register_pdf_font()
    page_w, page_h = 13.333 * inch, 7.5 * inch
    c = canvas.Canvas(str(OUT_PDF), pagesize=(page_w, page_h))
    for idx, spec in enumerate(SLIDES, 1):
        draw_rect(c, 0, 0, page_w, page_h, WHITE, WHITE)
        draw_rect(c, 0, page_h - 0.17 * inch, page_w, 0.17 * inch, GREEN, GREEN)
        if spec.get("layout") == "title":
            draw_text(c, 0.7 * inch, 6.45 * inch, 3.5 * inch, "AI TECH REVIEW | 2026-04-28", 10, True, GREEN)
            draw_text(c, 0.7 * inch, 5.65 * inch, 11.2 * inch, spec["title"], 25, True, INK)
            draw_text(c, 0.72 * inch, 4.92 * inch, 10 * inch, spec["subtitle"], 16, True, DARK_GREEN)
            draw_text(c, 0.75 * inch, 4.25 * inch, 10 * inch, "\n".join(spec["body"]), 11, False, INK)
            draw_rect(c, 0.72 * inch, 1.65 * inch, 11.9 * inch, 0.75 * inch, PALE, GRAY)
            draw_text(c, 0.95 * inch, 2.0 * inch, 11 * inch, "작업 규칙: 먼저 탐색하고, 선택 기준을 세운 뒤 실행한다.", 14, True, GREEN)
        else:
            draw_text(c, 0.42 * inch, 6.88 * inch, 9.2 * inch, spec["title"], 17, True, INK)
            draw_text(c, 0.44 * inch, 6.42 * inch, 9.9 * inch, spec.get("subtitle", ""), 9.8, True, DARK_GREEN)
            draw_text(c, 12.2 * inch, 6.88 * inch, 0.7 * inch, f"{idx:02d}", 10, True, GREEN)
            if "cards" in spec:
                cards = spec["cards"]
                cols = 2 if len(cards) in (4, 3) else 3
                card_w = (5.9 if cols == 2 else 3.8) * inch
                start_x = 0.68 * inch
                start_y = 4.72 * inch
                for i, (title, body) in enumerate(cards):
                    row = i // cols
                    col = i % cols
                    x = start_x + col * (card_w + 0.35 * inch)
                    y = start_y - row * 1.62 * inch
                    h = (1.32 if len(cards) > 3 else 1.45) * inch
                    accent = [GREEN, BLUE, ORANGE, RED][i % 4]
                    draw_rect(c, x, y, card_w, h, WHITE, GRAY)
                    draw_rect(c, x, y, 0.055 * inch, h, accent, accent)
                    draw_text(c, x + 0.16 * inch, y + h - 0.32 * inch, card_w - 0.3 * inch, title, 10, True, accent)
                    draw_text(c, x + 0.16 * inch, y + h - 0.67 * inch, card_w - 0.3 * inch, body, 8.6, False, INK)
            if "table" in spec:
                rows = spec["table"]
                x, y, w, h = 0.65 * inch, 1.2 * inch, 12.0 * inch, 4.75 * inch
                col_w = [0.25 * w, 0.49 * w, 0.26 * w]
                row_h = h / len(rows)
                for r, row in enumerate(rows):
                    cx = x
                    cy = y + h - (r + 1) * row_h
                    fill = PALE if r % 2 else WHITE
                    for cc, value in enumerate(row):
                        draw_rect(c, cx, cy, col_w[cc], row_h, fill, GRAY)
                        draw_text(c, cx + 0.08 * inch, cy + row_h - 0.27 * inch, col_w[cc] - 0.16 * inch, value, 8.8, cc == 0, GREEN if cc == 0 else INK)
                        cx += col_w[cc]
            if "bullets" in spec:
                draw_text(c, 0.9 * inch, 3.35 * inch, 10.6 * inch, "\n".join(f"- {b}" for b in spec["bullets"]), 10.2, False, INK)
            if "flow" in spec:
                items = spec["flow"]
                n = len(items)
                x0 = 0.72 * inch
                y = 1.55 * inch if "bullets" in spec else 0.95 * inch
                w = 11.8 * inch / n
                for i, item in enumerate(items):
                    draw_rect(c, x0 + i * w, y, w - 0.13 * inch, 0.55 * inch, LIGHT_GREEN if i % 2 == 0 else PALE, GREEN)
                    draw_text(c, x0 + i * w + 0.08 * inch, y + 0.31 * inch, w - 0.3 * inch, item, 8.3, True, DARK_GREEN)
        draw_text(c, 0.45 * inch, 0.27 * inch, 12.2 * inch, spec["footer"], 6.8, False, MUTED)
        c.showPage()
    c.save()


def main():
    render_pptx()
    render_pdf()
    print(OUT_PPTX)
    print(OUT_PDF)


if __name__ == "__main__":
    main()
